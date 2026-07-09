from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──
BG_DARK    = RGBColor(0x12, 0x12, 0x1A)
BG_CARD    = RGBColor(0x1E, 0x1E, 0x2E)
BG_CARD2   = RGBColor(0x16, 0x16, 0x22)
GREEN      = RGBColor(0x00, 0xC8, 0x53)
GREEN_DIM  = RGBColor(0x00, 0x80, 0x35)
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xAA, 0xAA, 0xBB)
GRAY_LIGHT = RGBColor(0xCC, 0xCC, 0xDD)
GRAY_DIM   = RGBColor(0x66, 0x66, 0x77)
CYAN       = RGBColor(0x00, 0xBF, 0xFF)
RED_SOFT   = RGBColor(0xFF, 0x45, 0x45)
YELLOW     = RGBColor(0xFF, 0xD6, 0x00)

# ── Helpers ──
def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_line(slide, x1, y1, x2, y2, color=GRAY_DIM, width=1.5):
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape

def add_icon_circle(slide, left, top, size, icon_text, bg_color, text_color=WHITE):
    circle = add_circle(slide, left, top, size, bg_color)
    txBox = add_text(slide, left, top, size, size, icon_text, size=int(size*20), color=text_color, bold=True, align=PP_ALIGN.CENTER)
    txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    txBox.text_frame.word_wrap = False
    return circle

def add_bullet_text(slide, left, top, width, height, items, size=14, color=GRAY_LIGHT, bullet_color=GREEN):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "▸ "
        run.font.size = Pt(size)
        run.font.color.rgb = bullet_color
        run.font.bold = True
        run.font.name = 'Segoe UI'
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(size)
        run2.font.color.rgb = color
        run2.font.name = 'Segoe UI'

# ═══════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

# Top accent line
add_rect(slide, 0, 0, 13.333, 0.06, GREEN)

# University name
add_text(slide, 0.8, 0.4, 12, 0.5,
         "UNIVERSIDAD NACIONAL HERMILIO VALDIZÁN", 13, GRAY, bold=False)
add_text(slide, 0.8, 0.7, 12, 0.4,
         "Facultad de Ingeniería Industrial, de Sistemas y Mecatrónica", 11, GRAY_DIM)

# Main title
add_text(slide, 0.8, 1.8, 11.5, 1.2,
         "Sistema de Riego Inteligente\npara un Huerto Universitario", 44, WHITE, bold=True)

# Subtitle
add_text(slide, 0.8, 3.3, 11.5, 0.8,
         "Gemelo Digital y Simulación de Infraestructura de Red", 22, GREEN, bold=False)

# Divider
add_rect(slide, 0.8, 4.3, 3, 0.04, GREEN)

# Bottom info
add_text(slide, 0.8, 4.6, 6, 0.4,
         "Integrantes:", 11, GRAY, bold=True)
integrantes = [
    "Orbezo Falcon, Abel Isaac",
    "Bernardo Cuellar, Jean Carlos",
    "Gilian Aguilar, Phatrick Sahit",
    "Dias Valles, Russell",
    "De la Cruz Guillen, Julio Cesar"
]
add_bullet_text(slide, 0.8, 4.9, 6, 2, integrantes, size=12, color=GRAY_LIGHT, bullet_color=GREEN_DIM)

add_text(slide, 0.8, 6.6, 6, 0.4,
         "Docente: Dra. Ines Eusebia Jesus Tolentino", 11, GRAY_DIM)
add_text(slide, 0.8, 6.9, 6, 0.4,
         "Ciencias Sociales y del Comportamiento Humano  |  Huánuco – Perú", 10, GRAY_DIM)

# Right side decorative
add_circle(slide, 10.5, 1.5, 2.5, RGBColor(0x00, 0xC8, 0x53))
add_circle(slide, 10.8, 1.8, 1.9, RGBColor(0x00, 0x80, 0x35))
add_text(slide, 10.5, 1.5, 2.5, 2.5, "🌿", size=60, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 2 — PROBLEMA
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, ORANGE)

add_text(slide, 0.8, 0.4, 12, 0.5,
         "EL PROBLEMA", 12, ORANGE, bold=True)
add_text(slide, 0.8, 0.8, 12, 0.8,
         "Crisis Alimentaria y Brecha Tecnológica", 36, WHITE, bold=True)

# Left card
add_rect(slide, 0.8, 1.9, 5.8, 4.8, BG_CARD, GRAY_DIM)
add_text(slide, 1.1, 2.1, 5.2, 0.5,
         "🌍  Contexto Global", 18, CYAN, bold=True)
add_bullet_text(slide, 1.1, 2.7, 5.2, 3.5, [
    "La FAO (2022) señala deficiencias en infraestructura para alimentos frescos",
    "La ONU (2021) exige modelos de producción que minimicen impacto ambiental",
    "El agua y el suelo son recursos vitales y finitos",
    "Brecha entre formación teórica y resolución de problemas reales"
], size=13, color=GRAY_LIGHT, bullet_color=ORANGE)

# Right card
add_rect(slide, 6.8, 1.9, 5.8, 4.8, BG_CARD, GRAY_DIM)
add_text(slide, 7.1, 2.1, 5.2, 0.5,
         "🏫  Contexto Local (UNHEVAL)", 18, ORANGE, bold=True)
add_bullet_text(slide, 7.1, 2.7, 5.2, 3.5, [
    "Ausencia de huerto inteligente en la universidad",
    "Toma de decisiones empírica sin sensores",
    "Carencia de infraestructura tecnológica en áreas de cultivo",
    "Falta de espacios para aprendizaje práctico"
], size=13, color=GRAY_LIGHT, bullet_color=ORANGE)

# Bottom highlight
add_rect(slide, 0.8, 6.0, 11.8, 0.8, RGBColor(0x33, 0x22, 0x00), ORANGE)
add_text(slide, 1.1, 6.1, 11.2, 0.6,
         "▸ Problema central: Ausencia de un modelo de huerto universitario inteligente que articule producción agroecológica con tecnología", 14, YELLOW, bold=False)

# ═══════════════════════════════════════════════
# SLIDE 3 — OBJETIVOS
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, GREEN)

add_text(slide, 0.8, 0.4, 12, 0.5,
         "OBJETIVOS", 12, GREEN, bold=True)
add_text(slide, 0.8, 0.8, 12, 0.8,
         "¿Qué queremos lograr?", 36, WHITE, bold=True)

# General
add_rect(slide, 0.8, 1.8, 11.8, 1.5, BG_CARD, GREEN)
add_text(slide, 1.1, 1.9, 11.2, 0.4,
         "OBJETIVO GENERAL", 13, GREEN, bold=True)
add_text(slide, 1.1, 2.3, 11.2, 0.9,
         "Diseñar y planificar un modelo de huerto universitario inteligente, técnica y presupuestalmente viable, que articule producción agroecológica, innovación tecnológica, aprendizaje práctico y sostenibilidad en la comunidad universitaria de la UNHEVAL.",
         14, GRAY_LIGHT)

# Especificos
add_text(slide, 0.8, 3.6, 12, 0.5,
         "OBJETIVOS ESPECÍFICOS", 13, GREEN, bold=True)

objs = [
    ("01", "Diseñar la propuesta técnica del huerto, considerando organización del espacio y módulos de cultivo"),
    ("02", "Planificar la incorporación progresiva de tecnologías IoT (sensores, monitoreo)"),
    ("03", "Estructurar propuesta presupuestal y operativa para implementación gradual"),
    ("04", "Diseñar acciones formativas con aprendizaje experiencial y educación ambiental"),
    ("05", "Establecer un modelo replicable que contribuya a la sostenibilidad institucional")
]

for i, (num, txt) in enumerate(objs):
    y = 4.1 + i * 0.62
    add_circle(slide, 1.0, y, 0.42, GREEN)
    add_text(slide, 1.0, y, 0.42, 0.42, num, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.6, y + 0.02, 10.8, 0.5, txt, 14, GRAY_LIGHT)

# ═══════════════════════════════════════════════
# SLIDE 4 — DESCRIPCIÓN DEL SISTEMA
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, CYAN)

add_text(slide, 0.8, 0.4, 12, 0.5,
         "DESCRIPCIÓN DEL SISTEMA", 12, CYAN, bold=True)
add_text(slide, 0.8, 0.8, 12, 0.8,
         "¿Cómo funciona el Huerto Inteligente?", 36, WHITE, bold=True)

# Architecture diagram (simplified)
# ESP32 box
add_rect(slide, 0.8, 1.9, 3.2, 2.2, BG_CARD, CYAN)
add_text(slide, 0.8, 2.0, 3.2, 0.5, "⚡  ESP32", 20, CYAN, bold=True, align=PP_ALIGN.CENTER)
add_bullet_text(slide, 1.0, 2.5, 2.8, 1.5, [
    "Lee 10 sensores cada 0.5s",
    "Controla relays (bomba, ventilador, pulverizador)",
    "Lazo cerrado automático"
], size=11, color=GRAY_LIGHT, bullet_color=CYAN)

# Arrow 1
add_text(slide, 4.0, 2.7, 0.8, 0.5, "→", size=30, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# Bridge box
add_rect(slide, 4.8, 1.9, 3.2, 2.2, BG_CARD, GREEN)
add_text(slide, 4.8, 2.0, 3.2, 0.5, "🌐  Bridge", 20, GREEN, bold=True, align=PP_ALIGN.CENTER)
add_bullet_text(slide, 5.0, 2.5, 2.8, 1.5, [
    "API en FastAPI (Render)",
    "Piggyback: sensores + comandos",
    "Endpoints REST"
], size=11, color=GRAY_LIGHT, bullet_color=GREEN)

# Arrow 2
add_text(slide, 8.0, 2.7, 0.8, 0.5, "→", size=30, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# Supabase box
add_rect(slide, 8.8, 1.9, 3.8, 2.2, BG_CARD, ORANGE)
add_text(slide, 8.8, 2.0, 3.8, 0.5, "🗄️  Supabase", 20, ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_text(slide, 9.0, 2.5, 3.4, 1.5, [
    "PostgreSQL en la nube",
    "6 tablas: sensores, lecturas, alertas...",
    "Realtime para el frontend"
], size=11, color=GRAY_LIGHT, bullet_color=ORANGE)

# Frontend box
add_rect(slide, 2.5, 4.5, 8.5, 2.3, BG_CARD, YELLOW)
add_text(slide, 2.5, 4.6, 8.5, 0.5, "🖥️  Frontend Web (SPA Brutalista)", 20, YELLOW, bold=True, align=PP_ALIGN.CENTER)
add_bullet_text(slide, 2.8, 5.1, 8, 1.5, [
    "Dashboard en tiempo real: gráfico, sensores, actuadores, configuración de umbrales",
    "Simulación de sequía, pH, temperatura, humedad → envía alertas al ESP32",
    "OverrideManager: ciclo visual de 3 fases (DEGRADING → HOLDING → RECOVERING)",
    "Juego educativo interactivo para aprender sobre riego inteligente"
], size=12, color=GRAY_LIGHT, bullet_color=YELLOW)

# Arrow down
add_text(slide, 6.2, 4.1, 1, 0.5, "↓", size=30, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 5 — HARDWARE POR MÓDULO
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, GREEN)

add_text(slide, 0.8, 0.4, 12, 0.5,
         "HARDWARE POR MÓDULO", 12, GREEN, bold=True)
add_text(slide, 0.8, 0.8, 12, 0.8,
         "Cada invernadero = 1 ESP32 + 10 sensores + actuadores", 28, WHITE, bold=True)

# Cards
components = [
    ("⚡", "ESP32", "Microcontrolador\nWi-Fi + Bluetooth\nDoble núcleo", CYAN),
    ("🌡️", "DHT22", "Temperatura y\nHumedad ambiente\n1 por invernadero", ORANGE),
    ("💧", "Humedad Suelo", "4 sensores capacitivos\n1 por maceta\nLectura analógica", GREEN),
    ("🧪", "pH Suelo", "4 sensores\n1 por maceta\nRango 0-14", RED_SOFT),
    ("📺", "OLED 0.96\"", "Display I2C\n5 páginas\nEstado en tiempo real", GRAY_LIGHT),
    ("🔀", "MUX 74HC4067", "16 canales analógicos\nMultiplexa 8 sensores\n1 pin ADC del ESP32", CYAN),
    ("💡", "Shift Register", "74HC595\n8 LEDs (4 verde + 4 rojo)\n3 GPIO del ESP32", GREEN),
    ("🔌", "Relays", "3 por maceta\nBomba / Ventilador /\nPulverizador", ORANGE),
]

for i, (icon, name, desc, color) in enumerate(components):
    col = i % 4
    row = i // 4
    x = 0.8 + col * 3.1
    y = 1.9 + row * 2.6
    add_rect(slide, x, y, 2.8, 2.2, BG_CARD, color)
    add_text(slide, x, y + 0.1, 2.8, 0.5, icon, size=28, align=PP_ALIGN.CENTER, color=color)
    add_text(slide, x, y + 0.6, 2.8, 0.4, name, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.15, y + 1.05, 2.5, 1.1, desc, size=11, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 6 — PIN MAPPING / CONEXIONES
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, CYAN)

add_text(slide, 0.8, 0.4, 12, 0.5,
         "CONECIONES Y PIN MAPPING", 12, CYAN, bold=True)
add_text(slide, 0.8, 0.8, 12, 0.8,
         "Distribución de pines del ESP32", 32, WHITE, bold=True)

# Left: MUX
add_rect(slide, 0.8, 1.8, 5.5, 4.5, BG_CARD, CYAN)
add_text(slide, 1.0, 1.9, 5.1, 0.5,
         "MUX CD74HC4067 — Sensores Analógicos", 16, CYAN, bold=True)
mux_items = [
    "S0=GPIO27  |  S1=GPIO16  |  S2=GPIO17  |  S3=GPIO18",
    "SIG=GPIO34 (ADC del ESP32)",
    "",
    "C0 → Humedad Suelo MAC-1    C4 → Humedad Suelo MAC-3",
    "C1 → pH MAC-1                      C5 → pH MAC-3",
    "C2 → Humedad Suelo MAC-2    C6 → Humedad Suelo MAC-4",
    "C3 → pH MAC-2                      C7 → pH MAC-4",
    "",
    "DHT22 = GPIO4 (compartido)"
]
add_bullet_text(slide, 1.1, 2.5, 5.1, 3.5, mux_items, size=12, color=GRAY_LIGHT, bullet_color=CYAN)

# Right: Relays + LEDs
add_rect(slide, 6.6, 1.8, 6.0, 4.5, BG_CARD, GREEN)
add_text(slide, 6.8, 1.9, 5.6, 0.5,
         "Actuadores y LEDs", 16, GREEN, bold=True)

relay_items = [
    "MAC-1: Bomba=GPIO32 | Vent=GPIO14 | Pulv=GPIO5",
    "MAC-2: Bomba=GPIO23 | Vent=GPIO0  | Pulv=GPIO15",
    "MAC-3: Bomba=GPIO12 | Vent=GPIO11 | Pulv=GPIO10",
    "MAC-4: Bomba=GPIO8  | Vent=GPIO7  | Pulv=GPIO6",
    "",
    "74HC595 Shift Register:",
    "  DS=GPIO2 | SHCP=GPIO3 | STCP=GPIO19",
    "  Q0-Q3 = Verde (relay activo)",
    "  Q4-Q7 = Rojo (condición crítica)",
    "",
    "LEDs globales: Naranja=33 | Amarillo=25 | Verde=1"
]
add_bullet_text(slide, 6.9, 2.5, 5.5, 3.5, relay_items, size=11, color=GRAY_LIGHT, bullet_color=GREEN)

# Bottom note
add_rect(slide, 0.8, 6.5, 11.8, 0.6, RGBColor(0x11, 0x22, 0x11), GREEN_DIM)
add_text(slide, 1.1, 6.55, 11.2, 0.5,
         "▸ 50 sensores totales: 5 invernaderos × (2 compartidos + 4×2 por maceta) = 50 IDs en Supabase",
         13, GREEN, bold=False)

# ═══════════════════════════════════════════════
# SLIDE 7 — FLUJO DE DATOS
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, GREEN)

add_text(slide, 0.8, 0.4, 12, 0.5,
         "FLUJO DE DATOS", 12, GREEN, bold=True)
add_text(slide, 0.8, 0.8, 12, 0.8,
         "Piggyback: todo en una sola llamada HTTP", 32, WHITE, bold=True)

# Step 1
add_rect(slide, 0.5, 1.8, 2.5, 4.5, BG_CARD, CYAN)
add_circle(slide, 1.3, 2.0, 0.6, CYAN)
add_text(slide, 1.3, 2.0, 0.6, 0.6, "1", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 0.5, 2.7, 2.5, 0.5, "ESP32 Lee", 16, CYAN, bold=True, align=PP_ALIGN.CENTER)
add_bullet_text(slide, 0.7, 3.2, 2.1, 2.5, [
    "DHT22: temp + hum_amb",
    "MUX: 4× hum_suelo",
    "MUX: 4× pH",
    "Total: 10 sensores",
    "Cada 0.5 segundos"
], size=11, color=GRAY_LIGHT, bullet_color=CYAN)

# Arrow
add_text(slide, 3.0, 3.5, 0.5, 0.5, "→", size=28, color=GREEN, bold=True)

# Step 2
add_rect(slide, 3.5, 1.8, 2.5, 4.5, BG_CARD, GREEN)
add_circle(slide, 4.3, 2.0, 0.6, GREEN)
add_text(slide, 4.3, 2.0, 0.6, 0.6, "2", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 3.5, 2.7, 2.5, 0.5, "Envía al Bridge", 16, GREEN, bold=True, align=PP_ALIGN.CENTER)
add_bullet_text(slide, 3.7, 3.2, 2.1, 2.5, [
    "POST /api/sync",
    "Device: 1",
    "Sensors: {0:25.3,...}",
    "Pending: true",
    "Cada 7 segundos"
], size=11, color=GRAY_LIGHT, bullet_color=GREEN)

# Arrow
add_text(slide, 6.0, 3.5, 0.5, 0.5, "→", size=28, color=GREEN, bold=True)

# Step 3
add_rect(slide, 6.5, 1.8, 2.5, 4.5, BG_CARD, ORANGE)
add_circle(slide, 7.3, 2.0, 0.6, ORANGE)
add_text(slide, 7.3, 2.0, 0.6, 0.6, "3", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 6.5, 2.7, 2.5, 0.5, "Bridge Procesa", 16, ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_text(slide, 6.7, 3.2, 2.1, 2.5, [
    "Batch INSERT 10 filas",
    "Consulta commands",
    "Consulta overrides",
    "Limpia alertas viejas",
    "Todo en paralelo"
], size=11, color=GRAY_LIGHT, bullet_color=ORANGE)

# Arrow
add_text(slide, 9.0, 3.5, 0.5, 0.5, "→", size=28, color=GREEN, bold=True)

# Step 4
add_rect(slide, 9.5, 1.8, 3.3, 4.5, BG_CARD, YELLOW)
add_circle(slide, 10.8, 2.0, 0.6, YELLOW)
add_text(slide, 10.8, 2.0, 0.6, 0.6, "4", size=22, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 9.5, 2.7, 3.3, 0.5, "Respuesta", 16, YELLOW, bold=True, align=PP_ALIGN.CENTER)
add_bullet_text(slide, 9.7, 3.2, 2.9, 2.5, [
    "ok: true",
    "sensors_written: 10",
    "commands: [relays]",
    "overrides: [alertas]",
    "ESP32 aplica todo"
], size=11, color=GRAY_LIGHT, bullet_color=YELLOW)

# Bottom
add_rect(slide, 0.5, 6.5, 12.3, 0.7, RGBColor(0x11, 0x22, 0x11), GREEN_DIM)
add_text(slide, 0.8, 6.55, 11.7, 0.6,
         "▸ 1 llamada HTTP reemplaza 4: sensores + comandos + overrides + cleanup = bajo uso de RAM en ESP32",
         13, GREEN, bold=False)

# ═══════════════════════════════════════════════
# SLIDE 8 — SIMULACIÓN Y OVERRIDE
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, ORANGE)

add_text(slide, 0.8, 0.4, 12, 0.5,
         "SIMULACIÓN Y CICLO DE OVERRIDE", 12, ORANGE, bold=True)
add_text(slide, 0.8, 0.8, 12, 0.8,
         "Creando escenarios críticos para probar el lazo cerrado", 28, WHITE, bold=True)

# Left: buttons
add_rect(slide, 0.8, 1.8, 5.5, 2.8, BG_CARD, ORANGE)
add_text(slide, 1.0, 1.9, 5.1, 0.5,
         "Botones de Simulación", 16, ORANGE, bold=True)
sim_items = [
    "SEQUÍA → hum_suelo baja a 10% (per-maceta)",
    "pH CRÍTICO → pH sube a 10.0 (per-maceta)",
    "TEMP EXTREMA → temp sube a 45°C (compartido)",
    "HUMEDAD BAJA → hum_amb baja a 15% (compartido)",
    "",
    "Validación: si valor dentro de umbrales → RECHAZA"
]
add_bullet_text(slide, 1.1, 2.4, 5.1, 2.2, sim_items, size=12, color=GRAY_LIGHT, bullet_color=ORANGE)

# Right: cycle
add_rect(slide, 6.6, 1.8, 6.0, 2.8, BG_CARD, RED_SOFT)
add_text(slide, 6.8, 1.9, 5.6, 0.5,
         "Ciclo de Override (ESP32)", 16, RED_SOFT, bold=True)
cycle_items = [
    "1. DEGRADING (21s): valor actual → crítico",
    "    7 pasos graduales × 3s cada uno",
    "2. HOLDING (2s): se mantiene en crítico",
    "3. RECOVERING (60s): crítico → setpoint ideal",
    "    20 pasos graduales × 3s cada uno",
    "TTL: 120 segundos → expira automáticamente"
]
add_bullet_text(slide, 6.9, 2.4, 5.5, 2.2, cycle_items, size=12, color=GRAY_LIGHT, bullet_color=RED_SOFT)

# Bottom: timeline
add_rect(slide, 0.8, 4.9, 11.8, 2.2, BG_CARD, GRAY_DIM)
add_text(slide, 1.0, 5.0, 11.4, 0.5,
         "Timeline de Visualización en Frontend (x10 más lento)", 16, YELLOW, bold=True)

phases = [
    ("WAITING", "30s", GRAY_DIM, 1.5),
    ("DEGRADING", "210s", RED_SOFT, 4.0),
    ("HOLDING", "20s", ORANGE, 1.2),
    ("RECOVERING", "100s", GREEN, 3.5),
    ("DONE", "", GRAY_DIM, 0.8)
]

x_start = 1.2
for name, dur, color, width in phases:
    add_rect(slide, x_start, 5.7, width, 0.7, color, None)
    add_text(slide, x_start, 5.7, width, 0.35, name, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if dur:
        add_text(slide, x_start, 6.05, width, 0.35, dur, size=10, color=WHITE, align=PP_ALIGN.CENTER)
    x_start += width + 0.15

add_text(slide, 1.0, 6.6, 11.4, 0.4,
         "▸ Actuadores visuales (bomba, ventilador, pulverizador) solo se muestran ON durante RECOVERING",
         12, GRAY, bold=False)

# ═══════════════════════════════════════════════
# SLIDE 9 — LAZO CERRADO
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, GREEN)

add_text(slide, 0.8, 0.4, 12, 0.5,
         "LAZO CERRADO AUTOMÁTICO", 12, GREEN, bold=True)
add_text(slide, 0.8, 0.8, 12, 0.8,
         "El ESP32 toma decisiones sin intervención humana", 30, WHITE, bold=True)

# 3 cards
actuadores = [
    ("💧", "BOMBA DE AGUA", "Humedad del suelo", [
        "ON si hum_suelo < 30%",
        "OFF si hum_suelo > 55%",
        "Riego por goteo individual"
    ], CYAN),
    ("🌀", "VENTILADOR", "Temperatura y humedad", [
        "ON si temp > 35°C",
        "ON si hum_amb > 85%",
        "OFF si temp < 30°C y hum_amb < 75%"
    ], GREEN),
    ("💨", "PULVERIZADOR", "Humedad ambiente", [
        "ON si hum_amb < 20%",
        "ON si temp > 35°C",
        "OFF si hum_amb > 50%"
    ], ORANGE),
]

for i, (icon, name, sensor, rules, color) in enumerate(actuadores):
    x = 0.8 + i * 4.1
    add_rect(slide, x, 1.8, 3.8, 4.5, BG_CARD, color)
    add_text(slide, x, 1.9, 3.8, 0.5, icon, size=36, align=PP_ALIGN.CENTER, color=color)
    add_text(slide, x, 2.5, 3.8, 0.5, name, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, 3.0, 3.8, 0.4, f"Controla: {sensor}", size=12, color=GRAY, align=PP_ALIGN.CENTER)
    add_rect(slide, x + 0.2, 3.5, 3.4, 0.03, color)
    add_bullet_text(slide, x + 0.3, 3.7, 3.2, 2.5, rules, size=13, color=GRAY_LIGHT, bullet_color=color)

# Bottom
add_rect(slide, 0.8, 6.5, 11.8, 0.7, RGBColor(0x11, 0x22, 0x11), GREEN_DIM)
add_text(slide, 1.1, 6.55, 11.2, 0.5,
         "▸ Evaluación continua cada 0.5s → LED naranja (crítico), amarillo (alerta), verde (normal)",
         13, GREEN, bold=False)

# ═══════════════════════════════════════════════
# SLIDE 10 — BASE DE DATOS
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, CYAN)

add_text(slide, 0.8, 0.4, 12, 0.5,
         "BASE DE DATOS — SUPABASE", 12, CYAN, bold=True)
add_text(slide, 0.8, 0.8, 12, 0.8,
         "PostgreSQL en la nube con 6 tablas principales", 30, WHITE, bold=True)

tables = [
    ("monitoreo_lecturas", "Lecturas de sensores", "sensor_id, valor_lectura, fecha_hora\n50 sensores, IDs 0-49\nBatch insert (10 filas/POST)", CYAN),
    ("sensores", "Catálogo de sensores", "50 registros\nTipo, maceta, fórmula de ID\nDispositivo asociado", GREEN),
    ("simulacion_alertas", "Alertas de simulación", "tipo, sensor, valor, maceta\nactiva=true → override activo\nCleanup 300s", ORANGE),
    ("control_actuadores", "Comandos de relays", "dispositivo, actuador, pin\nPENDIENTE → EJECUTADO\nVia piggyback", YELLOW),
    ("configuracion", "Umbrales por dispositivo", "JSONB: temp, humAmb,\nhumSuelo, ph (min/max)\nUpsert on save", RED_SOFT),
    ("dispositivos", "Nodos ESP32", "5 registros (IDs 1-5)\n1 por invernadero\nEstado de conexión", GRAY_LIGHT),
]

for i, (name, subtitle, desc, color) in enumerate(tables):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.8 + row * 2.6
    add_rect(slide, x, y, 3.8, 2.2, BG_CARD, color)
    add_text(slide, x + 0.15, y + 0.1, 3.5, 0.4, name, size=14, color=color, bold=True)
    add_text(slide, x + 0.15, y + 0.45, 3.5, 0.3, subtitle, size=11, color=GRAY)
    add_text(slide, x + 0.15, y + 0.85, 3.5, 1.2, desc, size=11, color=GRAY_LIGHT)

# ═══════════════════════════════════════════════
# SLIDE 11 — FRONTEND
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, YELLOW)

add_text(slide, 0.8, 0.4, 12, 0.5,
         "FRONTEND — PANEL DE CONTROL", 12, YELLOW, bold=True)
add_text(slide, 0.8, 0.8, 12, 0.8,
         "SPA Brutalista con diseño industrial", 32, WHITE, bold=True)

# Features
features = [
    ("📊", "Dashboard", "Gráfico en tiempo real\ncada 2 segundos\nDual Y-axis", GREEN),
    ("⚙️", "Configuración", "Umbrales guardados\nen Supabase\n4 parámetros", CYAN),
    ("🎮", "Simulación", "4 botones críticos\nValidación umbrales\nOverride visual", ORANGE),
    ("📋", "Tabla Sensores", "50 sensores en vivo\nFiltrado por módulo\nActualización auto", YELLOW),
    ("🎮", "Juego Educativo", "Aprender riego\n8 amenazas\nCombo streaks", RED_SOFT),
]

for i, (icon, name, desc, color) in enumerate(features):
    x = 0.5 + i * 2.5
    add_rect(slide, x, 1.8, 2.3, 3.5, BG_CARD, color)
    add_text(slide, x, 1.9, 2.3, 0.5, icon, size=32, align=PP_ALIGN.CENTER, color=color)
    add_text(slide, x, 2.5, 2.3, 0.5, name, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, x + 0.3, 3.0, 1.7, 0.03, color)
    add_text(slide, x + 0.15, 3.2, 2.0, 1.8, desc, size=11, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

# Tech stack
add_rect(slide, 0.8, 5.6, 11.8, 1.5, BG_CARD, GRAY_DIM)
add_text(slide, 1.0, 5.7, 11.4, 0.4,
         "Stack del Frontend", 14, GRAY_LIGHT, bold=True)
techs = "HTML5  •  CSS3 (Brutalist)  •  JavaScript (SPA)  •  Chart.js  •  ModuleLoader (?v=7 cache bust)  •  EventBus  •  Hash Router  •  JetBrains Mono"
add_text(slide, 1.0, 6.1, 11.4, 0.8, techs, 12, GRAY)

# ═══════════════════════════════════════════════
# SLIDE 12 — ODS
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, GREEN)

add_text(slide, 0.8, 0.4, 12, 0.5,
         "OBJETIVOS DE DESARROLLO SOSTENIBLE (ODS)", 12, GREEN, bold=True)
add_text(slide, 0.8, 0.8, 12, 0.8,
         "Contribución del proyecto a la Agenda 2030", 32, WHITE, bold=True)

ods = [
    ("2", "Hambre Cero", "Producción de alimentos frescos\nen el campus universitario", RGBColor(0xDD, 0xA6, 0x3A)),
    ("4", "Educación de Calidad", "Aprendizaje experiencial\ncon IoT y agricultura", RGBColor(0xC5, 0x19, 0x2D)),
    ("6", "Agua Limpia", "Riego inteligente que\noptimiza el uso del agua", RGBColor(0x26, 0xBD, 0xE2)),
    ("9", "Industria e Innovación", "Tecnología IoT aplicada\na producción agrícola", RGBColor(0xFD, 0x6A, 0x38)),
    ("11", "Ciudades Sostenibles", "Huerto urbano replicable\nen comunidades", RGBColor(0xFD, 0x9D, 0x24)),
    ("12", "Producción Responsable", "Uso eficiente de recursos\norgánicos y locales", RGBColor(0xBF, 0x8B, 0x2E)),
    ("13", "Acción Climática", "Reducción de desperdicio\nhídrico y energético", RGBColor(0x3F, 0x7E, 0x44)),
    ("17", "Alianzas", "Proyecto interdisciplinario\nuniversidad-comunidad", RGBColor(0x19, 0x48, 0x6A)),
]

for i, (num, name, desc, color) in enumerate(ods):
    col = i % 4
    row = i // 4
    x = 0.5 + col * 3.2
    y = 1.8 + row * 2.6
    add_rect(slide, x, y, 2.9, 2.2, color, None)
    add_text(slide, x, y + 0.1, 2.9, 0.5, f"ODS {num}", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + 0.6, 2.9, 0.4, name, size=14, color=WHITE, bold=False, align=PP_ALIGN.CENTER)
    add_rect(slide, x + 0.3, y + 1.1, 2.3, 0.03, WHITE)
    add_text(slide, x + 0.15, y + 1.25, 2.6, 0.8, desc, size=11, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 13 — PRESUPUESTO
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, ORANGE)

add_text(slide, 0.8, 0.4, 12, 0.5,
         "PRESUPUESTO", 12, ORANGE, bold=True)
add_text(slide, 0.8, 0.8, 12, 0.8,
         "Inversión total: S/ 1,545 (electrónica)", 32, WHITE, bold=True)

# Table header
add_rect(slide, 0.8, 1.8, 11.8, 0.5, GREEN_DIM)
add_text(slide, 1.0, 1.85, 4, 0.4, "Componente", 12, WHITE, bold=True)
add_text(slide, 5.0, 1.85, 3, 0.4, "Modelo", 12, WHITE, bold=True)
add_text(slide, 8.0, 1.85, 1.5, 0.4, "Cant.", 12, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 9.5, 1.85, 1.5, 0.4, "P.Unit.", 12, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 11.0, 1.85, 1.5, 0.4, "Subtotal", 12, WHITE, bold=True, align=PP_ALIGN.CENTER)

budget = [
    ("Microcontrolador", "ESP32 DevKit", "5", "S/ 35", "S/ 175"),
    ("Sensor Humedad", "Capacitivo v1.2/v2.0", "20", "S/ 10", "S/ 200"),
    ("Sensor pH", "pH-sensor Breakout", "20", "S/ 25", "S/ 500"),
    ("Sensor Ambiental", "DHT22", "5", "S/ 10", "S/ 50"),
    ("Pantalla Local", "OLED 0.96\" I2C", "5", "S/ 18", "S/ 90"),
    ("Expansor Salidas", "74HC595", "5", "S/ 3", "S/ 15"),
    ("Módulo Relé", "Relé 1-4 canales", "10", "S/ 12", "S/ 120"),
    ("Bomba de Agua", "Sumergible 5V", "5", "S/ 8", "S/ 40"),
    ("Ventilador", "Mini DC 5-12V", "5", "S/ 10", "S/ 50"),
    ("Panel Solar + Bat.", "Kit 5-6V + 18650", "5", "S/ 35", "S/ 175"),
    ("LEDs Indicadores", "Rojo/Verde", "20", "S/ 1", "S/ 20"),
    ("Accesorios", "Cables, MUX, etc.", "-", "-", "S/ 110"),
]

for i, (comp, model, cant, punit, subtotal) in enumerate(budget):
    y = 2.35 + i * 0.35
    bg = BG_CARD if i % 2 == 0 else BG_CARD2
    add_rect(slide, 0.8, y, 11.8, 0.35, bg)
    add_text(slide, 1.0, y + 0.02, 4, 0.3, comp, 11, GRAY_LIGHT)
    add_text(slide, 5.0, y + 0.02, 3, 0.3, model, 11, GRAY)
    add_text(slide, 8.0, y + 0.02, 1.5, 0.3, cant, 11, GRAY_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, 9.5, y + 0.02, 1.5, 0.3, punit, 11, GRAY_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, 11.0, y + 0.02, 1.5, 0.3, subtotal, 11, GRAY_LIGHT, align=PP_ALIGN.CENTER)

# Total
y_total = 2.35 + len(budget) * 0.35 + 0.1
add_rect(slide, 0.8, y_total, 11.8, 0.5, GREEN_DIM)
add_text(slide, 1.0, y_total + 0.05, 8, 0.4, "TOTAL ELECTRÓNICOS", 14, WHITE, bold=True)
add_text(slide, 11.0, y_total + 0.05, 1.5, 0.4, "S/ 1,545", 14, WHITE, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 14 — CONCLUSIONES
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, GREEN)

add_text(slide, 0.8, 0.4, 12, 0.5,
         "CONCLUSIONES", 12, GREEN, bold=True)
add_text(slide, 0.8, 0.8, 12, 0.8,
         "Resultados clave del proyecto", 36, WHITE, bold=True)

conclusions = [
    ("01", "La integración IoT en agricultura universitaria es factible con recursos modestos. Un sistema funcional por menos de S/ 500 por módulo."),
    ("02", "El equipo desarrolló competencias prácticas en arquitectura IoT, redes, gestión de proyectos y presupuestación técnica."),
    ("03", "La metodología combina trabajo de campo, investigación, simulación en Wokwi y Cisco Packet Tracer."),
    ("04", "El huerto sienta bases para innovación tecnológica y sostenibilidad en la UNHEVAL y la región Huánuco."),
]

for i, (num, txt) in enumerate(conclusions):
    y = 1.9 + i * 1.3
    add_rect(slide, 0.8, y, 11.8, 1.1, BG_CARD, GREEN)
    add_circle(slide, 1.1, y + 0.2, 0.6, GREEN)
    add_text(slide, 1.1, y + 0.2, 0.6, 0.6, num, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.9, y + 0.2, 10.4, 0.7, txt, 15, GRAY_LIGHT)

# Bottom
add_rect(slide, 0.8, 6.3, 11.8, 0.8, BG_CARD, CYAN)
add_text(slide, 1.0, 6.4, 11.4, 0.6,
         "▸ Repositorio: github.com/pgilian13-collab/huerto-puente  |  Demo: huerto-puente.onrender.com",
         13, CYAN, bold=False, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 15 — CIERRE
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.06, GREEN)

add_text(slide, 0, 2.5, 13.333, 1,
         "Gracias", 60, WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rect(slide, 5.5, 3.5, 2.3, 0.04, GREEN)

add_text(slide, 0, 3.8, 13.333, 0.6,
         "Sistema de Riego Inteligente — Huerto Universitario UNHEVAL", 18, GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 0, 4.5, 13.333, 0.5,
         "Integrantes: Orbezo · Bernardo · Gilian · Dias · De la Cruz", 14, GRAY_DIM, align=PP_ALIGN.CENTER)

add_text(slide, 0, 5.2, 13.333, 0.5,
         "Dra. Ines Eusebia Jesus Tolentino  |  Ciencias Sociales y del Comportamiento Humano", 12, GRAY_DIM, align=PP_ALIGN.CENTER)

add_circle(slide, 5.9, 5.8, 1.5, GREEN)
add_circle(slide, 6.1, 6.0, 1.1, RGBColor(0x00, 0x80, 0x35))
add_text(slide, 5.9, 5.8, 1.5, 1.5, "🌿", size=40, align=PP_ALIGN.CENTER)

# ── Save ──
output_path = os.path.join(os.getcwd(), "Huerto_Inteligente_UNHEVAL.pptx")
prs.save(output_path)
print(f"Presentación guardada en: {output_path}")
print(f"Total de diapositivas: {len(prs.slides)}")
